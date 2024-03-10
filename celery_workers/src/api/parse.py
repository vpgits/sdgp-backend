import logging
import re
import os
from io import BytesIO
from supabase import Client
from pypdf import PdfReader
from supabase import Client
from docx import Document
from pptx import Presentation
import json


from celery_workers.src.api.database import add_pages_to_supabase


# no test
def parse_pages(path: str, supabase_client: Client, document_id: str) -> None:
    # this function will parse an entire pdf into an array of pages
    try:
        download_file(
            supabase_client, path, document_id
        )  # this function will download a file from supabase storage
        pages = file_type(f"{path}.{get_file_type(supabase_client, document_id)}")
        # this function will return the type of file
        add_pages_to_supabase(
            supabase_client, pages, document_id
        )  # this function will add the pages to supabase
    except FileNotFoundError:
        logging.error("File not found")
    except Exception as e:
        logging.error(str(e))
        raise e


# test
def file_type(path: str) -> str:
    # this function will return the type of file
    try:
        if path.endswith(".pdf"):
            return read_pdf(path)
        elif path.endswith(".docx"):
            return read_docx(path)
        elif path.endswith(".pptx"):
            return read_pptx(path)
        else:
            logging.error("File type not supported: " + path)
            raise Exception(
                "File type not supported"
            )  # Fix: Replace 'throw' with 'raise' and add a newline
    except Exception as e:
        logging.error("Error while getting file type: " + str(e))
        raise e


# no test
def does_pages_exist(supabase: Client, document_id: str) -> bool:
    try:
        data = (
            supabase.from_("documents").select("data").eq("id", document_id).execute()
        )
        if not data.data[0].get("data"):
            return False
        else:
            return True
    except Exception as e:
        logging.error(str(e))
        raise e


def read_pdf(path: str) -> list[str]:
    # Parse the PDF file and return a list of text from each page
    pages = []
    try:
        with open(f"./resources/{path}", "rb") as file:
            reader = PdfReader(BytesIO(file.read()))
            for page in reader.pages:
                cleaned_page = re.sub(r"\s+", " ", page.extract_text())
                pages.append(cleaned_page)
        return pages
    except FileNotFoundError as e:
        logging.error(f"File not found: {e}")
        raise
    except Exception as e:
        logging.error(f"Error parsing PDF: {str(e)}")
        raise


def read_docx(path: str) -> list[str]:
    # Parse the DOCX file and return a list of text from each page
    pages = []
    try:
        doc = Document(f"./resources/{path}")
        for para in doc.paragraphs:
            pages.append(para.text)
        return pages
    except FileNotFoundError as e:
        logging.error(f"File not found: {e}")
        raise
    except Exception as e:
        logging.error(f"Error parsing DOCX: {str(e)}")
        raise


def read_pptx(path: str) -> list[str]:
    # Parse the PPTX file and return a list of text from each page
    pages = []
    try:
        prs = Presentation(f"./resources/{path}")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pages.append(shape.text)
        return pages
    except FileNotFoundError as e:
        logging.error(f"File not found: {e}")
        raise
    except Exception as e:
        logging.error(f"Error parsing PPTX: {str(e)}")
        raise


# test
def get_file_type(supabase_client: Client, document_id: str) -> str:
    data = (
        supabase_client.table("documents")
        .select("file_type")
        .eq("id", document_id)
        .execute()
    )
    return data.data[0].get("file_type")


# test
def download_file(supabase_client: Client, path: str, document_id: str) -> None:
    # this function will download a file from supabase storage
    file_type = get_file_type(supabase_client, document_id)
    file_path = f"{path}.{file_type}"
    res = supabase_client.storage.from_("public/pdf").download(file_path)
    try:
        # if the ./resources folder does not exist, create it
        os.makedirs(os.path.dirname(f"./resources/{file_path}"), exist_ok=True)
        with open(f"./resources/{file_path}", "wb+") as f:
            f.write(res)
    except Exception as e:
        logging.error("Error Downloading file: " + str(e))
        raise e


from celery_workers.src.api.requests import create_document_summary


# no test needed
def create_document_summary_context(
    pages: list, supabase_client: Client, document_id: str
):
    try:
        response = create_document_summary(pages)
        supabase_client.table("documents").update({"summary": json.loads(response)}).eq(
            "id", document_id
        ).execute()
    except Exception as e:
        logging.error("Error while creating document context: " + str(e))
        raise e
